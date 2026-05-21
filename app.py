# ============================================================
#  📘 STUDY BUDDY — AI Learning Assistant (v3)
# ============================================================
#  Run in Google Colab. Auto-installs packages and launches a
#  Gradio app with a public share link.
#
#  KEY CHANGES IN v3:
#   - Uses OpenAI API (instead of Groq) via a single ai_chat() wrapper
#   - RAG (TF-IDF + cosine, sklearn) used ONLY for:
#         * Ask Question From Notes
#         * Active Recall Practice
#     (with a keyword fallback if sklearn is missing)
#   - Topic-specific generation (PDF + topic = only that topic)
#   - PDF content cleaning (removes admin/title-page noise)
#   - Study plan = clean TABLE
#   - Quiz = clean TABLE (no raw JSON shown to the user)
#   - Active recall = FRESH questions every session
#   - Weak questions = no duplicates (deduped by concept/text)
#   - Clean, human-readable downloads
#   - Clickable/collapsible mindmap + graphical dashboard
#   - Q&A understands English / Roman Urdu / Urdu
#   - Friendly error handling (no Gradio output mismatch crashes)
# ============================================================

# ============================================================
#  🔑 STEP 0 — PUT YOUR OPENAI API KEY HERE (once)
# ============================================================
#  Get a key at: https://platform.openai.com/api-keys
#  Either paste it below, OR set environment variable OPENAI_API_KEY.
#  In Colab you can also do: import os; os.environ["OPENAI_API_KEY"]="sk-..."
API_KEY = "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"          # <-- paste your OpenAI key, e.g. "sk-proj-..."
MODEL = "gpt-4o-mini" # balanced: fast + cheap. Change to "gpt-4o" if you want.

# ---------- Install dependencies (safe to re-run) ----------
import subprocess, sys

def _pip_install(pkg):
    """Install a package quietly; ignore failures so the app still runs."""
    try:
        subprocess.run([sys.executable, "-m", "pip", "install", "-q", pkg],
                       check=True)
    except Exception as e:
        print(f"(warning) could not install {pkg}: {e}")

for _p in ["openai", "gradio", "pypdf", "fpdf2", "python-pptx",
           "qrcode[pil]", "scikit-learn", "matplotlib"]:
    _pip_install(_p)

# ---------- Imports ----------
import os
import io
import json
import re
import html
import hashlib
import random
import tempfile
from datetime import datetime

from pypdf import PdfReader
from fpdf import FPDF
import gradio as gr

# Optional libraries (handled gracefully if missing)
try:
    import qrcode
    QR_AVAILABLE = True
except Exception:
    QR_AVAILABLE = False

try:
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except Exception:
    OPENAI_AVAILABLE = False

# sklearn for RAG (TF-IDF); fallback to keyword search if missing
try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
    SKLEARN_AVAILABLE = True
except Exception:
    SKLEARN_AVAILABLE = False

# matplotlib for progress charts (handled gracefully if missing)
try:
    import matplotlib
    matplotlib.use("Agg")   # no display needed (Colab/server safe)
    import matplotlib.pyplot as plt
    MPL_AVAILABLE = True
except Exception:
    MPL_AVAILABLE = False


# ============================================================
#  GLOBAL CONFIG & MEMORY
# ============================================================
PROGRESS_FILE = "study_progress.json"
USERS_FILE = "users.json"          # stores accounts (email, name, hashed password)

# Per-student session memory (RAM, this run only):
#   raw_notes, clean_notes, focused_notes, selected_topic,
#   pdf_title, sources, summary, quiz, plan_rows, tree,
#   recall_quiz, recall_index, chunks (for RAG)
session_memory = {}

_ai_client = None

def _get_client():
    """Create the OpenAI client from API_KEY var or OPENAI_API_KEY env var."""
    global _ai_client
    if _ai_client is not None:
        return _ai_client
    if not OPENAI_AVAILABLE:
        return None
    key = API_KEY.strip() or os.environ.get("OPENAI_API_KEY", "").strip()
    if not key:
        return None
    try:
        _ai_client = OpenAI(api_key=key)
        return _ai_client
    except Exception:
        return None


# ============================================================
#  STEP 1 (kept) — UNIFIED AI WRAPPER  ai_chat()
# ============================================================
def ai_chat(system_prompt, user_prompt, temperature=0.4):
    """
    Single entry point for all AI calls (now OpenAI).
    Never crashes — returns a friendly message if not configured.
    """
    client = _get_client()
    if client is None:
        return ("⚠ AI is not configured. Please set your OpenAI API key in the "
                "API_KEY variable at the top of the code (or the OPENAI_API_KEY "
                "environment variable). Get one at platform.openai.com/api-keys")
    try:
        resp = client.chat.completions.create(
            model=MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=temperature,
        )
        return resp.choices[0].message.content
    except Exception as e:
        return f"⚠ AI request failed: {str(e)}"


def _ai_failed(text):
    """True if ai_chat returned an error/placeholder, not real content."""
    if not text:
        return True
    return ("⚠ AI is not configured" in text) or ("⚠ AI request failed" in text)


# ============================================================
#  USER ACCOUNTS / AUTHENTICATION
# ============================================================
#  Accounts are stored in users.json as:
#    { "email@example.com": {"name": "...", "pw": "<sha256 hash>"} }
#  Passwords are hashed (sha256), never stored in plain text.
#  The user's EMAIL becomes their student_id, so each account keeps
#  its own uploads, progress, weak questions, etc.

def _load_users():
    """Read all accounts from users.json (returns {} if none)."""
    try:
        if os.path.exists(USERS_FILE):
            with open(USERS_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}

def _save_users(data):
    """Write accounts back to users.json."""
    try:
        with open(USERS_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
    except Exception as e:
        print(f"(warning) could not save users: {e}")

def _hash_pw(password):
    """Turn a password into a secure sha256 hash."""
    return hashlib.sha256((password or "").encode("utf-8")).hexdigest()

def _valid_email(email):
    """Very simple email format check."""
    return bool(re.match(r"^[^@\s]+@[^@\s]+\.[^@\s]+$", (email or "").strip()))


def sign_up(email, full_name, password, confirm):
    """
    Create a new account.
    Returns (success_bool, message).
    """
    email = (email or "").strip().lower()
    full_name = (full_name or "").strip()
    if not _valid_email(email):
        return False, "❌ Please enter a valid email address."
    if not full_name:
        return False, "❌ Please enter your full name."
    if len(password or "") < 4:
        return False, "❌ Password must be at least 4 characters."
    if password != confirm:
        return False, "❌ Passwords do not match."

    users = _load_users()
    if email in users:
        return False, "❌ This email is already registered. Please sign in."
    users[email] = {"name": full_name, "pw": _hash_pw(password)}
    _save_users(users)
    return True, f"✅ Account created for {full_name}! You can now sign in."


def sign_in(email, password):
    """
    Check login details.
    Returns (success_bool, message, name).
    """
    email = (email or "").strip().lower()
    users = _load_users()
    if email not in users:
        return False, "❌ No account found. Please sign up first.", ""
    if users[email]["pw"] != _hash_pw(password):
        return False, "❌ Wrong password. Please try again.", ""
    return True, f"✅ Welcome back, {users[email]['name']}!", users[email]["name"]


# ============================================================
#  PROGRESS STORAGE (local JSON file)
# ============================================================
def load_progress():
    try:
        if os.path.exists(PROGRESS_FILE):
            with open(PROGRESS_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}

def save_progress(data):
    try:
        with open(PROGRESS_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
    except Exception as e:
        print(f"(warning) could not save progress: {e}")

def get_student_record(student_id):
    """Get (or create) a student's long-term progress record."""
    data = load_progress()
    if student_id not in data:
        data[student_id] = {
            "pdfs_studied": 0,
            "quizzes_generated": 0,
            "questions_attempted": 0,
            "correct": 0,
            "wrong": 0,
            "last_topic": "",
            "weak_topics": {},        # concept -> wrong count
            "weak_questions": {},     # qid -> record (deduped)
            "last_study": "",
            "sources": [],
        }
        save_progress(data)
    if isinstance(data[student_id].get("weak_questions"), list):
        data[student_id]["weak_questions"] = {}
    return data, data[student_id]


# ============================================================
#  PDF EXTRACTION
# ============================================================
def extract_text_from_pdf(file_obj):
    """Extract raw text from one PDF file object/path."""
    try:
        if hasattr(file_obj, "name"):
            reader = PdfReader(file_obj.name)
        else:
            reader = PdfReader(io.BytesIO(file_obj))
        text = ""
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
        return text.strip() if text.strip() else ""
    except Exception as e:
        return f"PDF Error: {str(e)}"


# ============================================================
#  STEP 2 (kept) — CONTENT CLEANING
# ============================================================
def rule_based_clean(raw_text):
    """Cheap first pass: drop obvious title-page / header / footer lines."""
    if not raw_text:
        return ""
    noise_keywords = [
        "university", "department", "professor", "lecturer", "instructor",
        "course code", "course title", "roll no", "roll number", "student name",
        "assignment", "submitted to", "submitted by", "semester", "session",
        "copyright", "all rights reserved", "table of contents", "faculty of",
        "date:", "name:", "id:", "page ",
    ]
    out = []
    for line in raw_text.splitlines():
        s = line.strip()
        if not s:
            continue
        low = s.lower()
        if re.fullmatch(r"(page\s*)?\d{1,4}", low):   # pure page numbers
            continue
        if any(k in low for k in noise_keywords) and len(s) < 80:
            continue
        out.append(s)
    return "\n".join(out).strip()


def clean_learning_content(raw_text):
    """Rule-based clean + AI smart clean. Falls back gracefully."""
    if not raw_text or "PDF Error" in raw_text:
        return raw_text or "No text extracted from PDF"
    pre = rule_based_clean(raw_text) or raw_text
    if _get_client() is None:
        return pre
    system = (
        "You are a strict academic content filter. From the raw PDF text, KEEP "
        "only real educational content: concepts, definitions, explanations, "
        "examples, theories, formulas, processes, comparisons, case studies, and "
        "important points. REMOVE university/department/professor/student names, "
        "course code, roll number, title page, dates, semester, page numbers, "
        "headers, footers, table of contents, copyright, repeated slide titles. "
        "Output ONLY the cleaned learning content as plain text."
    )
    result = ai_chat(system, f"Clean this text:\n\n{pre[:12000]}", temperature=0.1)
    if _ai_failed(result) or len(result) < 30:
        return pre
    return result.strip()


# ============================================================
#  STEP 1 (prompt) — TOPIC-SPECIFIC EXTRACTION
# ============================================================
def extract_topic_content(clean_notes, topic):
    """
    Return only the parts of clean_notes about `topic`.
    Returns (focused_text, found_boolean).
    """
    if not topic or not topic.strip():
        return clean_notes, True
    if not clean_notes:
        return "Topic not found in the uploaded PDF.", False

    topic = topic.strip()
    keyword_hit = topic.lower() in clean_notes.lower()

    if _get_client() is not None:
        system = (
            "You are a precise study assistant. The user wants ONLY content about "
            f"the topic: '{topic}'. From the notes, extract and rewrite (in clear "
            "student-friendly language) ONLY the sections, definitions, and "
            "examples about this topic. If the topic is genuinely NOT covered, "
            "reply with exactly: TOPIC_NOT_FOUND . Otherwise output only the "
            "relevant rephrased content (do not copy verbatim)."
        )
        result = ai_chat(system, f"Topic: {topic}\n\nNotes:\n{clean_notes[:12000]}",
                         temperature=0.2)
        if not _ai_failed(result):
            if "TOPIC_NOT_FOUND" in result and not keyword_hit:
                return "Topic not found in the uploaded PDF.", False
            if "TOPIC_NOT_FOUND" not in result and len(result) > 30:
                return result.strip(), True

    if keyword_hit:   # fallback keyword slice
        lines = clean_notes.splitlines()
        keep, capture = [], False
        for ln in lines:
            if topic.lower() in ln.lower():
                capture = True
            if capture:
                keep.append(ln)
        return ("\n".join(keep).strip() or clean_notes), True

    return "Topic not found in the uploaded PDF.", False


# ============================================================
#  BUILD CLEAN + FOCUSED NOTES (multi-PDF + topic)
# ============================================================
def build_notes(pdf_files, topic_text):
    """Returns a dict with raw/clean/focused notes + status info."""
    sources, warnings, parts = [], [], []
    pdf_title = ""

    if pdf_files:
        if not isinstance(pdf_files, list):
            pdf_files = [pdf_files]
        for f in pdf_files:
            if f is None:
                continue
            name = os.path.basename(getattr(f, "name", "uploaded.pdf"))
            if not pdf_title:
                pdf_title = os.path.splitext(name)[0]
            raw = extract_text_from_pdf(f)
            if not raw or "PDF Error" in raw:
                warnings.append(f"Could not read text from {name}, skipped.")
                continue
            cleaned = clean_learning_content(raw)
            if cleaned and "No text" not in cleaned:
                parts.append(cleaned)
                sources.append(name)
            else:
                warnings.append(f"No usable content in {name}, skipped.")

    raw_notes = "\n\n".join(parts).strip()
    clean_notes = raw_notes
    topic = (topic_text or "").strip()
    focused_notes = clean_notes
    topic_found = True
    topic_msg = ""

    if clean_notes and topic:
        focused, found = extract_topic_content(clean_notes, topic)
        topic_found = found
        if found:
            focused_notes = focused
            topic_msg = f"✅ Topic '{topic}' found in PDF. Generating only this topic."
        else:
            focused_notes = ""
            topic_msg = ("❌ Topic not found in the uploaded PDF. "
                         "Please enter a topic that exists in the PDF.")
        if not pdf_title:
            pdf_title = topic
    elif clean_notes and not topic:
        focused_notes = clean_notes
        topic_msg = "✅ Using the full PDF content (no topic filter)."
    elif not clean_notes and topic:
        gen = ai_chat(
            "You are an expert teacher. Produce clear student-friendly study "
            "notes on the topic: definitions, key concepts, explanations, "
            "examples, and important points. Plain text.",
            f"Create study notes on the topic: {topic}", temperature=0.5)
        if not _ai_failed(gen):
            focused_notes = clean_notes = gen
            sources.append(f"Topic: {topic}")
            pdf_title = topic
            topic_msg = f"✅ Generated AI learning material for topic '{topic}'."
        else:
            topic_found = False
            topic_msg = "❌ Could not generate topic notes (AI not available)."
    else:
        topic_msg = "⚠ Please upload a PDF or enter a topic."

    return {
        "raw_notes": raw_notes, "clean_notes": clean_notes,
        "focused_notes": focused_notes, "selected_topic": topic,
        "pdf_title": pdf_title or "Study Material", "sources": sources,
        "warnings": warnings, "topic_found": topic_found, "topic_msg": topic_msg,
    }


# ============================================================
#  STEP 10 (prompt) — RAG: chunking + retrieval
#  (used ONLY for Q&A and Active Recall)
# ============================================================
def make_chunks(text, size=600, overlap=80):
    """Split text into overlapping word chunks for retrieval."""
    words = (text or "").split()
    chunks, i = [], 0
    while i < len(words):
        chunk = " ".join(words[i:i + size])
        if chunk.strip():
            chunks.append(chunk.strip())
        i += size - overlap
    return chunks or ([text] if text else [])


def retrieve_chunks(query, chunks, k=4):
    """
    Return the top-k most relevant chunks for a query.
    Uses TF-IDF + cosine similarity; falls back to keyword overlap.
    """
    if not chunks:
        return []
    if len(chunks) <= k:
        return chunks

    if SKLEARN_AVAILABLE:
        try:
            vec = TfidfVectorizer(stop_words="english")
            matrix = vec.fit_transform(chunks + [query])
            sims = cosine_similarity(matrix[-1], matrix[:-1]).flatten()
            top = sims.argsort()[::-1][:k]
            return [chunks[i] for i in top]
        except Exception:
            pass

    # Fallback: simple keyword overlap scoring
    q_words = set(re.findall(r"\w+", query.lower()))
    scored = []
    for c in chunks:
        c_words = set(re.findall(r"\w+", c.lower()))
        scored.append((len(q_words & c_words), c))
    scored.sort(key=lambda x: -x[0])
    return [c for _, c in scored[:k]]


def get_rag_context(student_id, query, k=4):
    """Build a focused context string using RAG over stored chunks."""
    m = _mem(student_id)
    if not m:
        return ""
    chunks = m.get("chunks")
    if not chunks:
        chunks = make_chunks(m.get("focused_notes", ""))
        m["chunks"] = chunks
    top = retrieve_chunks(query, chunks, k=k)
    return "\n\n".join(top)


# ============================================================
#  STEP 11 (prompt) — STUDENT-FRIENDLY AI MODULES
# ============================================================
def layered_summary(notes, topic=""):
    """Layered, student-friendly summary (rephrased, faithful to source)."""
    if not notes:
        return "No content available."
    focus = f" Focus ONLY on the topic: '{topic}'." if topic else ""
    system = (
        "You are a friendly teacher. Rephrase the source into EASY, "
        "student-friendly language (do not copy verbatim). Stay faithful to the "
        "source meaning." + focus +
        " Produce a LAYERED summary using these markdown sections:\n"
        "## 🔹 One-line Overview\n## 🔹 Key Concepts\n## 🔹 Detailed Explanation\n"
        "## 🔹 Important Definitions\n## 🔹 Exam-focused Points\n"
        "## 🔹 Common Mistakes\nUse bullets. Keep it clear, not robotic."
    )
    return ai_chat(system, f"Source content:\n{notes[:10000]}")


def _parse_json_array(raw):
    """Safely pull a JSON array out of an AI response."""
    try:
        cleaned = raw.replace("```json", "").replace("```", "").strip()
        s, e = cleaned.find("["), cleaned.rfind("]")
        if s != -1 and e != -1:
            cleaned = cleaned[s:e + 1]
        data = json.loads(cleaned)
        if isinstance(data, list):
            return data
    except Exception:
        pass
    return []


def generate_quiz_json(notes, topic="", n=6, avoid=None, temperature=0.5):
    """Varied quiz with rich fields (concept/scenario/application/etc)."""
    if not notes:
        return []
    focus = f" All questions must be about: '{topic}'." if topic else ""
    avoid_txt = ""
    if avoid:
        avoid_txt = (" Do NOT repeat or lightly reword these previous "
                     f"questions: {' | '.join(avoid[-12:])}.")
    system = (
        "You are a friendly exam question writer. Create a VARIED set of MCQs "
        "from the content, rephrased clearly (not copied)." + focus + avoid_txt +
        " Mix question_type: 'concept','definition','scenario','application',"
        "'comparison','example'. At least 2 must be 'scenario' (real situation). "
        "Keep questions short and clear. Respond with ONLY valid JSON, no "
        "markdown:\n"
        '[{"id":"q1","question":"...","options":{"A":"..","B":"..","C":"..",'
        '"D":".."},"correct_answer":"A","explanation":"short, student-friendly",'
        '"topic":"...","question_type":"scenario","difficulty":"easy|medium|hard",'
        '"scenario_context":"... or empty"}]'
    )
    raw = ai_chat(system, f"Make {n} MCQs from:\n{notes[:9000]}",
                  temperature=temperature)
    data = _parse_json_array(raw)
    for i, q in enumerate(data, 1):
        q.setdefault("id", f"q{i}")
        q.setdefault("question_type", "concept")
        q.setdefault("difficulty", "medium")
        q.setdefault("topic", topic or "General")
        if "correct_answer" not in q and "answer" in q:
            q["correct_answer"] = q["answer"]
    return data


def generate_fresh_recall_quiz(student_id, focused_notes, previous_questions, topic=""):
    """
    STEP 5 (prompt) — fresh recall set each session, using RAG context
    + randomization, avoiding old wording.
    """
    # Use a random slice of retrieved chunks as the context seed
    seed_query = topic or "key concepts definitions examples"
    context = get_rag_context(student_id, seed_query, k=5) or focused_notes
    # randomize so each session feels different
    random.seed(datetime.now().timestamp())
    hint = random.choice([
        "Make them scenario-heavy.", "Mix in application questions.",
        "Add comparison/difference questions.", "Use fresh real-world examples.",
    ])
    return generate_quiz_json(context + "\n\n" + hint, topic=topic, n=5,
                              avoid=previous_questions, temperature=0.95)


def generate_study_plan_rows(notes, topic="", days=7, minutes=60):
    """STEP 3 (prompt) — study plan as table rows."""
    if not notes:
        return []
    focus = f" Focus only on: '{topic}'." if topic else ""
    # Adjust length if topic is small/large (rough heuristic)
    suggested = 3 if len(notes) < 1500 else (5 if len(notes) < 4000 else days)
    system = (
        "You are a study coach. Build a personalized study plan." + focus +
        f" The student has ~{minutes} minutes/day. Use about {suggested} days "
        "(small topic = fewer days). Base it ONLY on real concepts, rephrased "
        "simply. Respond with ONLY valid JSON array, no markdown:\n"
        '[{"day":"Day 1","topic":"...","key_points":"- point - point",'
        '"practice":"...","output":"..."}]'
    )
    raw = ai_chat(system, f"Content:\n{notes[:8000]}", temperature=0.4)
    return _parse_json_array(raw)


def generate_flashcards(notes, topic=""):
    """Front/Back flashcards (rephrased, topic-specific)."""
    if not notes:
        return []
    focus = f" Only about: '{topic}'." if topic else ""
    system = (
        "Create study flashcards from the content, rephrased simply." + focus +
        " Respond with ONLY valid JSON, no markdown:\n"
        '[{"front":"question or key term","back":"clear answer/explanation"}]. '
        "Make 8 flashcards."
    )
    raw = ai_chat(system, f"Content:\n{notes[:8000]}", temperature=0.4)
    return _parse_json_array(raw)


def generate_mindmap_tree(notes, topic="", root_title="Study Material"):
    """STEP 8 (prompt) — nested tree for clickable mindmap."""
    if not notes:
        return {"name": root_title, "children": []}
    focus = f" Focus only on: '{topic}'." if topic else ""
    system = (
        "Build a 3-level mindmap from the content." + focus +
        " Level 1 = main topics, level 2 = subtopics, level 3 = short details "
        "(definitions/examples/key points), rephrased simply. Respond with ONLY "
        "valid JSON, no markdown:\n"
        '{"name":"ROOT","children":[{"name":"Main Topic","children":'
        '[{"name":"Subtopic","children":[{"name":"detail"}]}]}]}'
    )
    raw = ai_chat(system, f"Content:\n{notes[:8000]}", temperature=0.4)
    try:
        cleaned = raw.replace("```json", "").replace("```", "").strip()
        s, e = cleaned.find("{"), cleaned.rfind("}")
        if s != -1 and e != -1:
            cleaned = cleaned[s:e + 1]
        tree = json.loads(cleaned)
        tree["name"] = root_title
        return tree
    except Exception:
        return {"name": root_title, "children": []}


# ============================================================
#  STEP 9 (prompt) — Q&A with RAG + English/Roman Urdu/Urdu
# ============================================================
def answer_from_notes(question, student_id):
    """Answer using RAG-retrieved notes; mirror the user's language."""
    m = _mem(student_id)
    if not m or not m.get("focused_notes"):
        return "❌ Please first generate notes (upload a PDF or enter a topic)."

    context = get_rag_context(student_id, question, k=4)
    if not context:
        context = m.get("focused_notes", "")[:6000]

    system = (
        "You are a friendly teacher answering a student. Use ONLY the provided "
        "notes as the source of truth. The question may be in English, Roman "
        "Urdu (Urdu in English letters), or Urdu script — understand it and "
        "reply in the SAME language style. Rephrase in easy human language; do "
        "not copy verbatim. If the notes do NOT contain the answer, reply "
        "exactly: 'This answer is not available in the uploaded PDF/topic "
        "content.' If you add a small helpful example beyond the notes, put it "
        "under a clearly labelled line: 'Extra simple explanation:'. Never mix "
        "outside facts into the main answer without that label."
    )
    return ai_chat(system, f"Notes:\n{context}\n\nQuestion: {question}")


# ============================================================
#  PRESENTATION HELPERS (clean HTML / markdown for the UI)
# ============================================================
def _esc(s):
    return html.escape(str(s if s is not None else ""))


def quiz_to_table_md(quiz):
    """Interactive, styled quiz CARDS (no raw JSON). Answers hidden until clicked."""
    if not quiz:
        return ("<div style='padding:14px;color:#0f172a;font-family:system-ui'>"
                "📭 No quiz yet. Click <b>Generate Study Material</b> first.</div>")
    # type -> (icon, badge color)
    type_style = {
        "concept": ("📘", "#2563eb"), "definition": ("📖", "#7c3aed"),
        "scenario": ("🎬", "#db2777"), "application": ("🛠️", "#ea580c"),
        "comparison": ("⚖️", "#0891b2"), "example": ("💡", "#ca8a04"),
    }
    diff_style = {"easy": "#16a34a", "medium": "#d97706", "hard": "#dc2626"}
    cards = ""
    for i, q in enumerate(quiz, 1):
        opts = q.get("options", {})
        ans = str(q.get("correct_answer", q.get("answer", ""))).upper()
        qtype = q.get("question_type", "concept")
        icon, badge = type_style.get(qtype, ("📝", "#475569"))
        diff = q.get("difficulty", "medium")
        sc = q.get("scenario_context", "")

        scenario_html = ""
        if sc and str(sc).strip().lower() not in ("", "none", "empty"):
            scenario_html = (f"<div style='background:#fdf2f8;border-left:4px solid "
                             f"#db2777;padding:10px 12px;border-radius:8px;margin:8px 0;"
                             f"color:#0f172a;font-style:italic;'>🎬 <b>Scenario:</b> "
                             f"{_esc(sc)}</div>")

        opts_html = ""
        for k in ["A", "B", "C", "D"]:
            if k in opts:
                is_ans = (k == ans)
                if is_ans:
                    opts_html += (f"<div style='display:flex;gap:8px;padding:8px 12px;"
                                  f"margin:5px 0;border-radius:8px;background:#dcfce7;"
                                  f"border:1px solid #86efac;color:#0f172a;'>"
                                  f"<b style='color:#16a34a;'>{k}</b>"
                                  f"<span>{_esc(opts[k])}</span>"
                                  f"<span style='margin-left:auto;'>✅</span></div>")
                else:
                    opts_html += (f"<div style='display:flex;gap:8px;padding:8px 12px;"
                                  f"margin:5px 0;border-radius:8px;background:#f8fafc;"
                                  f"border:1px solid #e2e8f0;color:#0f172a;'>"
                                  f"<b style='color:#0f172a;'>{k}</b>"
                                  f"<span>{_esc(opts[k])}</span></div>")

        cards += f"""
        <div style="border:1px solid #e2e8f0;border-radius:16px;padding:18px;
                    margin-bottom:16px;background:#ffffff;box-shadow:0 2px 8px
                    rgba(0,0,0,0.05);color:#0f172a;">
          <div style="display:flex;gap:8px;flex-wrap:wrap;align-items:center;
                      margin-bottom:10px;">
            <span style="background:#0f172a;color:#fff;width:30px;height:30px;
                  display:flex;align-items:center;justify-content:center;
                  border-radius:50%;font-weight:700;">{i}</span>
            <span style="background:{badge};color:#fff;padding:3px 12px;
                  border-radius:999px;font-size:12px;font-weight:600;">
                  {icon} {_esc(qtype)}</span>
            <span style="background:{diff_style.get(diff,'#475569')};color:#fff;
                  padding:3px 12px;border-radius:999px;font-size:12px;">{_esc(diff)}</span>
            <span style="background:#fef9c3;color:#854d0e;padding:3px 12px;
                  border-radius:999px;font-size:12px;">🏷️ {_esc(q.get('topic','General'))}</span>
          </div>
          {scenario_html}
          <div style="font-weight:700;font-size:15px;margin:8px 0;color:#0f172a;">
            {_esc(q.get('question',''))}</div>
          {opts_html}
          <details style="margin-top:10px;">
            <summary style="cursor:pointer;color:#0891b2;font-weight:600;">
              💡 Show Explanation</summary>
            <div style="margin-top:6px;color:#0f172a;background:#ecfeff;padding:10px;
                 border-radius:8px;">{_esc(q.get('explanation',''))}</div>
          </details>
        </div>"""
    return (f"<div style='font-family:system-ui,Arial'>"
            f"<h3 style='color:#0f172a;'>📝 Quiz</h3>{cards}</div>")


def study_plan_to_md(rows, days, minutes):
    """STEP 3 (prompt) — study plan markdown table."""
    if not rows:
        return "_No study plan available. Generate study material first._"
    header = ("| Day | Topic Focus | Key Points | Practice Task | Output |\n"
              "|-----|-------------|------------|---------------|--------|\n")
    body = ""
    for r in rows:
        def cell(x):
            return str(r.get(x, "")).replace("\n", " ").replace("|", "/")
        body += (f"| {cell('day')} | {cell('topic')} | {cell('key_points')} "
                 f"| {cell('practice')} | {cell('output')} |\n")
    return (f"### 🗓 Study Plan ({minutes} min/day)\n\n" + header + body)


def build_mindmap_html(tree):
    """STEP 8 (prompt) — clickable/collapsible mindmap."""
    def node_html(node, level=0):
        name = _esc(node.get("name", ""))
        children = node.get("children", [])
        if level == 0:
            color, bg = "#38bdf8", "rgba(56,189,248,0.12)"
        elif level == 1:
            color, bg = "#fbbf24", "rgba(251,191,36,0.10)"
        else:
            color, bg = "#a7f3d0", "transparent"
        if children:
            inner = "".join(node_html(c, level + 1) for c in children)
            # All branches open by default so clicking a topic reveals its details.
            return (f"<details open style='margin-left:{level*16}px;'>"
                    f"<summary style='cursor:pointer;color:{color};"
                    f"font-weight:{700 if level<2 else 500};padding:6px 10px;"
                    f"background:{bg};border-radius:8px;margin:3px 0;"
                    f"list-style:none;user-select:none;'>"
                    f"{'🧠 ' if level==0 else ('📂 ' if level==1 else '• ')}{name}"
                    f"</summary>{inner}</details>")
        return (f"<div style='margin-left:{(level+1)*16}px;color:{color};"
                f"padding:3px 8px;'>• {name}</div>")
    try:
        body = node_html(tree, 0)
        return (f"<div style='font-family:system-ui,Arial;padding:16px;"
                f"background:#0f172a;border-radius:12px;line-height:1.7;'>"
                f"<div style='color:#cbd5e1;margin-bottom:10px;font-size:13px;'>"
                f"💡 All topics are open. Click any topic title to collapse or "
                f"expand it.</div>{body}</div>")
    except Exception:
        return "<pre>Mindmap unavailable.</pre>"


def mindmap_to_html_file(tree):
    """Downloadable mindmap HTML."""
    body = build_mindmap_html(tree)
    f = tempfile.NamedTemporaryFile(suffix=".html", delete=False,
                                    mode="w", encoding="utf-8")
    f.write("<!doctype html><html><head><meta charset='utf-8'>"
            "<title>Mindmap</title></head><body style='background:#0b1220'>"
            + body + "</body></html>")
    f.close()
    return f.name


# ============================================================
#  STEP 7 (prompt) — CLEAN DOWNLOADABLE FILES
# ============================================================
def _safe(s):
    """Make text safe for FPDF (latin-1)."""
    return str(s).encode("latin-1", "ignore").decode("latin-1")

def _md_to_plain(text):
    """Strip markdown symbols for clean human-readable files."""
    text = re.sub(r"#+\s*", "", text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    return text.strip()


def _styled_header(pdf, title, subtitle="", rgb=(37, 99, 235)):
    """Draw a colored title banner at the top of a PDF page."""
    r, g, b = rgb
    pdf.set_fill_color(r, g, b)
    pdf.rect(0, 0, 210, 26, "F")            # full-width colored band
    pdf.set_xy(10, 7)
    pdf.set_text_color(255, 255, 255)
    pdf.set_font("Arial", "B", 18)
    pdf.cell(0, 8, _safe(title), ln=1)
    if subtitle:
        pdf.set_x(10)
        pdf.set_font("Arial", "I", 10)
        pdf.cell(0, 6, _safe(subtitle), ln=1)
    pdf.set_text_color(0, 0, 0)             # back to black for body text
    pdf.ln(12)


def quiz_pdf_styled(quiz, title="Quiz"):
    """Beautiful, human-readable quiz PDF with colors & icons. TXT fallback."""
    if not quiz:
        quiz = []
    try:
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        _styled_header(pdf, "Study Buddy - Quiz", title, rgb=(37, 99, 235))
        type_icon = {"concept": "[Concept]", "definition": "[Definition]",
                     "scenario": "[Scenario]", "application": "[Application]",
                     "comparison": "[Comparison]", "example": "[Example]"}
        for i, q in enumerate(quiz, 1):
            opts = q.get("options", {})
            ans = str(q.get("correct_answer", q.get("answer", ""))).upper()
            qtype = q.get("question_type", "concept")
            diff = q.get("difficulty", "medium")
            # tag line
            pdf.set_fill_color(241, 245, 249)
            pdf.set_font("Arial", "B", 9)
            pdf.set_text_color(71, 85, 105)
            pdf.cell(0, 7, _safe(f"  Q{i}   {type_icon.get(qtype,'[MCQ]')}   "
                                 f"Difficulty: {diff}   Topic: {q.get('topic','General')}"),
                     ln=1, fill=True)
            pdf.set_text_color(0, 0, 0)
            # scenario
            sc = q.get("scenario_context", "")
            if sc and str(sc).strip().lower() not in ("", "none", "empty"):
                pdf.set_font("Arial", "I", 10)
                pdf.set_text_color(190, 24, 93)
                pdf.multi_cell(0, 6, _safe(f"Scenario: {sc}"))
                pdf.set_text_color(0, 0, 0)
            # question
            pdf.set_font("Arial", "B", 11)
            pdf.multi_cell(0, 7, _safe(q.get("question", "")))
            # options
            pdf.set_font("Arial", "", 11)
            for k in ["A", "B", "C", "D"]:
                if k in opts:
                    if k == ans:
                        pdf.set_text_color(22, 163, 74)   # green for correct
                        pdf.set_font("Arial", "B", 11)
                        pdf.multi_cell(0, 6, _safe(f"   {k}) {opts[k]}   (correct)"))
                        pdf.set_font("Arial", "", 11)
                        pdf.set_text_color(0, 0, 0)
                    else:
                        pdf.multi_cell(0, 6, _safe(f"   {k}) {opts[k]}"))
            # explanation
            pdf.set_text_color(8, 145, 178)
            pdf.set_font("Arial", "I", 10)
            pdf.multi_cell(0, 6, _safe(f"Explanation: {q.get('explanation','')}"))
            pdf.set_text_color(0, 0, 0)
            pdf.ln(4)
        out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        pdf.output(out.name)
        return out.name
    except Exception:
        return quiz_file_clean(quiz)   # plain-text fallback


def recall_pdf_styled(quiz, title="Active Recall"):
    """Styled PDF for a recall practice session."""
    return quiz_pdf_styled(quiz, title=title)


def flashcards_pdf(cards, reference_text="Study Buddy"):
    """Clean Front/Back flashcards PDF with optional QR code."""
    if not cards:
        cards = [{"front": "No flashcards generated", "back": "Try regenerating."}]
    try:
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Arial", "B", 18)
        pdf.cell(0, 14, _safe("Study Buddy  -  Flashcards"), ln=1, align="C")
        pdf.set_draw_color(180, 180, 180)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(6)
        if QR_AVAILABLE:
            try:
                img = qrcode.make((reference_text or "Study Buddy")[:120])
                qtmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                img.save(qtmp.name)
                pdf.image(qtmp.name, x=85, y=pdf.get_y(), w=35)
                pdf.ln(40)
                pdf.set_font("Arial", "I", 9)
                pdf.cell(0, 6, _safe("Scan for reference"), ln=1, align="C")
                pdf.ln(2)
            except Exception:
                pass
        for i, c in enumerate(cards, 1):
            pdf.set_fill_color(224, 242, 254)
            pdf.set_font("Arial", "B", 12)
            pdf.cell(0, 9, _safe(f"  Card {i}"), ln=1, fill=True)
            pdf.ln(1)
            pdf.set_font("Arial", "B", 11)
            pdf.multi_cell(0, 7, _safe(f"Front:  {c.get('front','')}"))
            pdf.set_font("Arial", "", 11)
            pdf.multi_cell(0, 7, _safe(f"Back:   {c.get('back','')}"))
            pdf.ln(4)
        out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        pdf.output(out.name)
        return out.name
    except Exception:
        f = tempfile.NamedTemporaryFile(suffix=".txt", delete=False,
                                        mode="w", encoding="utf-8")
        for i, c in enumerate(cards, 1):
            f.write(f"Card {i}\nFront: {c.get('front','')}\n"
                    f"Back: {c.get('back','')}\n\n")
        f.close()
        return f.name


def summary_pdf(summary_md, title="Summary"):
    """Styled summary PDF with colored header & section highlights. TXT fallback."""
    try:
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        _styled_header(pdf, "Study Buddy - Summary", title, rgb=(124, 58, 237))
        for raw in summary_md.splitlines():
            line = raw.rstrip()
            if not line.strip():
                continue
            # markdown heading -> colored bold heading
            if line.lstrip().startswith("#"):
                heading = _md_to_plain(line)
                pdf.ln(2)
                pdf.set_text_color(124, 58, 237)
                pdf.set_font("Arial", "B", 13)
                pdf.multi_cell(0, 8, _safe(heading))
                pdf.set_text_color(0, 0, 0)
            elif line.lstrip().startswith(("-", "*", "•")):
                pdf.set_font("Arial", "", 11)
                pdf.multi_cell(0, 7, _safe("  - " + _md_to_plain(line.lstrip("-*• "))))
            else:
                pdf.set_font("Arial", "", 11)
                pdf.multi_cell(0, 7, _safe(_md_to_plain(line)))
        out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        pdf.output(out.name)
        return out.name
    except Exception:
        f = tempfile.NamedTemporaryFile(suffix=".txt", delete=False,
                                        mode="w", encoding="utf-8")
        f.write(_md_to_plain(summary_md))
        f.close()
        return f.name


def study_plan_file(rows, days, minutes):
    """Styled study-plan PDF table. TXT fallback."""
    try:
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        _styled_header(pdf, "Study Buddy - Study Plan",
                       f"About {minutes} minutes per day", rgb=(5, 150, 105))
        if not rows:
            pdf.set_font("Arial", "", 11)
            pdf.cell(0, 8, _safe("No plan available."), ln=1)
        for idx, r in enumerate(rows):
            # alternating soft background per day block
            bg = (236, 253, 245) if idx % 2 == 0 else (240, 253, 250)
            pdf.set_fill_color(*bg)
            pdf.set_text_color(5, 150, 105)
            pdf.set_font("Arial", "B", 12)
            pdf.cell(0, 9, _safe(f"  {r.get('day','Day')}  -  {r.get('topic','')}"),
                     ln=1, fill=True)
            pdf.set_text_color(0, 0, 0)
            pdf.set_font("Arial", "", 11)
            pdf.multi_cell(0, 6, _safe(f"   Key Points : {r.get('key_points','')}"))
            pdf.multi_cell(0, 6, _safe(f"   Practice   : {r.get('practice','')}"))
            pdf.multi_cell(0, 6, _safe(f"   Output     : {r.get('output','')}"))
            pdf.ln(3)
        out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        pdf.output(out.name)
        return out.name
    except Exception:
        f = tempfile.NamedTemporaryFile(suffix=".txt", delete=False,
                                        mode="w", encoding="utf-8")
        f.write("STUDY BUDDY - STUDY PLAN\n")
        f.write(f"About {minutes} minutes per day\n" + "=" * 45 + "\n\n")
        for r in rows:
            f.write(f"{r.get('day','Day')}\n")
            f.write(f"  Topic Focus : {r.get('topic','')}\n")
            f.write(f"  Key Points  : {r.get('key_points','')}\n")
            f.write(f"  Practice    : {r.get('practice','')}\n")
            f.write(f"  Output      : {r.get('output','')}\n\n")
        f.close()
        return f.name


def quiz_file_clean(quiz):
    """Human-readable quiz TXT (not raw JSON)."""
    f = tempfile.NamedTemporaryFile(suffix=".txt", delete=False,
                                    mode="w", encoding="utf-8")
    f.write("STUDY BUDDY - QUIZ\n" + "=" * 40 + "\n\n")
    for i, q in enumerate(quiz, 1):
        f.write(f"Q{i}. [{q.get('question_type','concept')} | "
                f"{q.get('difficulty','medium')} | {q.get('topic','General')}]\n")
        sc = q.get("scenario_context", "")
        if sc and str(sc).strip().lower() not in ("", "none", "empty"):
            f.write(f"   Scenario: {sc}\n")
        f.write(f"   {q.get('question','')}\n")
        for k in ["A", "B", "C", "D"]:
            if k in q.get("options", {}):
                f.write(f"      {k}) {q['options'][k]}\n")
        f.write(f"   Correct: {q.get('correct_answer', q.get('answer',''))}\n")
        f.write(f"   Why: {q.get('explanation','')}\n\n")
    f.close()
    return f.name


def _build_progress_charts(rec):
    """Create chart images (pie for correct/wrong, bar for weak topics).
    Returns a list of temp PNG paths. Empty list if matplotlib missing."""
    if not MPL_AVAILABLE:
        return []
    paths = []
    correct = rec.get("correct", 0)
    wrong = rec.get("wrong", 0)
    try:
        # --- Pie chart: correct vs wrong ---
        if (correct + wrong) > 0:
            fig, ax = plt.subplots(figsize=(4, 3))
            ax.pie([correct, wrong], labels=["Correct", "Wrong"],
                   colors=["#16a34a", "#ef4444"], autopct="%1.0f%%",
                   startangle=90, textprops={"color": "#0f172a"})
            ax.set_title("Answer Accuracy", color="#0f172a")
            p = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            fig.savefig(p.name, bbox_inches="tight", dpi=110)
            plt.close(fig)
            paths.append(p.name)
        # --- Bar chart: weak topics ---
        weak = rec.get("weak_topics", {})
        if weak:
            items = sorted(weak.items(), key=lambda x: -x[1])[:6]
            labels = [t[:18] for t, _ in items]
            vals = [c for _, c in items]
            fig, ax = plt.subplots(figsize=(5, 3))
            ax.bar(labels, vals, color="#f59e0b")
            ax.set_title("Weak Topics (mistakes)", color="#0f172a")
            ax.tick_params(axis="x", rotation=30, labelcolor="#0f172a")
            ax.tick_params(axis="y", labelcolor="#0f172a")
            p = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            fig.savefig(p.name, bbox_inches="tight", dpi=110)
            plt.close(fig)
            paths.append(p.name)
    except Exception:
        pass
    return paths


def progress_report_file(student_id):
    """Styled progress-report PDF WITH charts (pie + bar). TXT fallback."""
    _, rec = get_student_record(student_id)
    attempted = rec.get("questions_attempted", 0)
    acc = (rec.get("correct", 0) / attempted * 100) if attempted else 0
    try:
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        _styled_header(pdf, "Study Buddy - Progress Report",
                       f"Student: {student_id}", rgb=(2, 132, 199))

        # stat lines (all black text)
        pdf.set_font("Arial", "", 12)
        stats = [
            ("PDFs / Topics studied", rec.get("pdfs_studied", 0)),
            ("Quizzes generated", rec.get("quizzes_generated", 0)),
            ("Questions attempted", attempted),
            ("Correct answers", rec.get("correct", 0)),
            ("Wrong answers", rec.get("wrong", 0)),
            ("Accuracy", f"{acc:.1f}%"),
            ("Last topic", rec.get("last_topic", "-")),
            ("Last study", rec.get("last_study", "-")),
        ]
        for label, val in stats:
            pdf.set_font("Arial", "B", 11)
            pdf.cell(60, 8, _safe(f"{label}:"), border=0)
            pdf.set_font("Arial", "", 11)
            pdf.cell(0, 8, _safe(str(val)), ln=1)
        pdf.ln(4)

        # charts
        charts = _build_progress_charts(rec)
        for cpath in charts:
            try:
                pdf.image(cpath, w=120)
                pdf.ln(4)
            except Exception:
                pass

        out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        pdf.output(out.name)
        return out.name
    except Exception:
        f = tempfile.NamedTemporaryFile(suffix=".txt", delete=False,
                                        mode="w", encoding="utf-8")
        f.write(f"STUDY BUDDY - PROGRESS REPORT ({student_id})\n" + "=" * 40 + "\n\n")
        f.write(f"PDFs/Topics studied : {rec.get('pdfs_studied',0)}\n")
        f.write(f"Quizzes generated   : {rec.get('quizzes_generated',0)}\n")
        f.write(f"Questions attempted : {attempted}\n")
        f.write(f"Correct answers     : {rec.get('correct',0)}\n")
        f.write(f"Wrong answers       : {rec.get('wrong',0)}\n")
        f.write(f"Accuracy            : {acc:.1f}%\n")
        f.close()
        return f.name


def summary_to_pptx(summary_md, title="Summary"):
    """Simple PPTX from summary; TXT fallback."""
    plain = _md_to_plain(summary_md)
    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation()
        bullets = [b.strip() for b in plain.splitlines() if b.strip()]
        layout = prs.slide_layouts[1]
        for i in range(0, len(bullets), 5):
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"{title} - Slide {i//5 + 1}"
            tf = slide.placeholders[1].text_frame
            for b in bullets[i:i + 5]:
                p = tf.add_paragraph()
                p.text = b
                p.font.size = Pt(16)
        out = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(out.name)
        return out.name
    except Exception:
        f = tempfile.NamedTemporaryFile(suffix=".txt", delete=False,
                                        mode="w", encoding="utf-8")
        f.write(plain)
        f.close()
        return f.name


# ============================================================
#  GRAPHICAL DASHBOARD
# ============================================================
def build_dashboard_html(student_id):
    """Graphical dashboard: stat cards + donut chart + weak-topic bar chart.
    All text is black for readability."""
    _, rec = get_student_record(student_id)
    attempted = rec.get("questions_attempted", 0)
    correct = rec.get("correct", 0)
    wrong = rec.get("wrong", 0)
    acc = (correct / attempted * 100) if attempted else 0

    def card(label, value, color):
        return (f"<div style='flex:1;min-width:110px;background:{color};"
                f"color:#fff;border-radius:14px;padding:14px;text-align:center;'>"
                f"<div style='font-size:26px;font-weight:800;'>{value}</div>"
                f"<div style='font-size:12px;opacity:.95;'>{label}</div></div>")

    cards = ("<div style='display:flex;gap:10px;flex-wrap:wrap;margin-bottom:16px;'>"
             + card("PDFs/Topics", rec.get("pdfs_studied", 0), "#0ea5e9")
             + card("Quizzes", rec.get("quizzes_generated", 0), "#8b5cf6")
             + card("Attempted", attempted, "#f59e0b")
             + card("Correct", correct, "#16a34a")
             + card("Wrong", wrong, "#ef4444") + "</div>")

    # --- Donut chart (CSS conic-gradient) for correct vs wrong ---
    if attempted > 0:
        donut = (
            f"<div style='display:flex;align-items:center;gap:18px;flex-wrap:wrap;"
            f"margin-bottom:16px;'>"
            f"<div style='width:130px;height:130px;border-radius:50%;background:"
            f"conic-gradient(#16a34a 0% {acc:.1f}%, #ef4444 {acc:.1f}% 100%);"
            f"display:flex;align-items:center;justify-content:center;'>"
            f"<div style='width:90px;height:90px;border-radius:50%;background:#fff;"
            f"display:flex;flex-direction:column;align-items:center;"
            f"justify-content:center;'>"
            f"<div style='font-size:22px;font-weight:800;color:#0f172a;'>{acc:.0f}%</div>"
            f"<div style='font-size:11px;color:#0f172a;'>accuracy</div></div></div>"
            f"<div style='color:#0f172a;font-size:14px;'>"
            f"<div>🟢 Correct: <b>{correct}</b></div>"
            f"<div>🔴 Wrong: <b>{wrong}</b></div>"
            f"<div>📝 Total: <b>{attempted}</b></div></div></div>")
    else:
        donut = ("<div style='color:#0f172a;padding:8px;font-size:14px;'>"
                 "Attempt some recall questions to see your performance chart.</div>")

    # --- Weak-topic horizontal bar chart ---
    weak = rec.get("weak_topics", {})
    if weak:
        items = sorted(weak.items(), key=lambda x: -x[1])
        maxc = max(c for _, c in items) or 1
        bars = ""
        for t, c in items:
            pct = (c / maxc) * 100
            bars += (
                f"<div style='margin-bottom:8px;color:#0f172a;'>"
                f"<div style='font-size:13px;margin-bottom:2px;'>{_esc(t)} "
                f"<b style='color:#ef4444;'>({c})</b></div>"
                f"<div style='background:#fee2e2;border-radius:8px;height:16px;'>"
                f"<div style='width:{pct:.0f}%;background:linear-gradient(90deg,"
                f"#f59e0b,#ef4444);height:100%;border-radius:8px;'></div></div></div>")
        weak_block = (f"<h4 style='margin:8px 0;color:#0f172a;'>⚠ Weak Topics</h4>{bars}")
    else:
        weak_block = ("<div style='color:#16a34a;padding:8px;font-weight:600;'>"
                      "🎉 No weak topics yet — great job!</div>")

    meta = (f"<div style='font-size:13px;color:#0f172a;margin-top:12px;'>"
            f"📖 Last topic: <b>{_esc(rec.get('last_topic','-'))}</b> &nbsp;|&nbsp; "
            f"🕒 Last study: <b>{_esc(rec.get('last_study','-'))}</b></div>")

    return (f"<div style='font-family:system-ui,Arial;padding:18px;"
            f"background:#f8fafc;border-radius:16px;color:#0f172a;'>"
            f"<h3 style='margin-top:0;color:#0f172a;'>📊 Progress Dashboard — "
            f"{_esc(student_id)}</h3>{cards}"
            f"<h4 style='margin:8px 0;color:#0f172a;'>🎯 Performance</h4>{donut}"
            f"{weak_block}{meta}</div>")


# ============================================================
#  SESSION HELPERS
# ============================================================
def _mem(student_id):
    return session_memory.get((student_id or "student123").strip())

def _focus_of(student_id):
    m = _mem(student_id)
    return (m or {}).get("focused_notes"), (m or {}).get("selected_topic", "")


# ============================================================
#  MAIN PIPELINE  (fixed output count — STEP 13 prompt)
# ============================================================
#  Outputs (11): summary, quiz_md, plan_md, mind_html,
#                flash, ppt, quizf, mindf, progressf,
#                status, dashboard
def _blank_outputs(student_id, msg):
    block = "<div style='padding:12px'>" + msg + "</div>"
    return (msg, msg, msg, block, None, None, None, None, None, msg,
            build_dashboard_html(student_id))

def full_pipeline(pdf_files, topic_text, student_id, days, minutes):
    student_id = (student_id or "student123").strip()
    days = int(days or 7)
    minutes = int(minutes or 60)

    if not pdf_files and not (topic_text and topic_text.strip()):
        return _blank_outputs(student_id, "⚠ Please upload a PDF or enter a topic.")

    info = build_notes(pdf_files, topic_text)

    if info["selected_topic"] and not info["topic_found"]:
        return _blank_outputs(student_id, info["topic_msg"])

    focused = info["focused_notes"]
    topic = info["selected_topic"]
    title = info["pdf_title"]

    if not focused:
        return _blank_outputs(student_id, "⚠ No usable content found.")

    # Save to session memory + build RAG chunks
    session_memory[student_id] = {
        "raw_notes": info["raw_notes"], "clean_notes": info["clean_notes"],
        "focused_notes": focused, "selected_topic": topic,
        "pdf_title": title, "sources": info["sources"],
        "chunks": make_chunks(focused),
    }

    try:
        summary = layered_summary(focused, topic)
    except Exception as e:
        summary = f"Summary failed: {e}"
    session_memory[student_id]["summary"] = summary

    try:
        quiz = generate_quiz_json(focused, topic=topic, n=6)
    except Exception:
        quiz = []
    session_memory[student_id]["quiz"] = quiz
    quiz_md = quiz_to_table_md(quiz)

    try:
        plan_rows = generate_study_plan_rows(focused, topic, days, minutes)
    except Exception:
        plan_rows = []
    session_memory[student_id]["plan_rows"] = plan_rows
    plan_md = study_plan_to_md(plan_rows, days, minutes)

    try:
        tree = generate_mindmap_tree(focused, topic, root_title=title)
    except Exception:
        tree = {"name": title, "children": []}
    session_memory[student_id]["tree"] = tree
    mind_html = build_mindmap_html(tree)

    try:
        cards = generate_flashcards(focused, topic)
        flash = flashcards_pdf(cards, reference_text=", ".join(info["sources"]))
    except Exception:
        flash = None
    try:
        ppt = summary_to_pptx(summary, title)
    except Exception:
        ppt = None
    try:
        quizf = quiz_pdf_styled(quiz, title=topic or title)
    except Exception:
        quizf = None
    try:
        mindf = mindmap_to_html_file(tree)
    except Exception:
        mindf = None

    try:
        data, rec = get_student_record(student_id)
        rec["pdfs_studied"] += len([s for s in info["sources"]
                                    if not s.startswith("Topic:")])
        rec["quizzes_generated"] += 1
        rec["last_topic"] = topic or (info["sources"][-1] if info["sources"] else title)
        rec["last_study"] = datetime.now().strftime("%Y-%m-%d %H:%M")
        for s in info["sources"]:
            if s not in rec["sources"]:
                rec["sources"].append(s)
        data[student_id] = rec
        save_progress(data)
    except Exception:
        pass

    status = info["topic_msg"] + "  |  📚 Sources: " + (", ".join(info["sources"]) or "—")
    if info["warnings"]:
        status += "  |  ⚠ " + " ".join(info["warnings"])

    return (summary, quiz_md, plan_md, mind_html,
            flash, ppt, quizf, mindf, progress_report_file(student_id),
            status, build_dashboard_html(student_id))


# ============================================================
#  RETRY / REGENERATE (focused content only)
# ============================================================
def retry_summary(student_id):
    notes, topic = _focus_of(student_id)
    if not notes:
        return "⚠ Generate study material first."
    s = layered_summary(notes, topic)
    _mem(student_id)["summary"] = s
    return s

def retry_quiz(student_id):
    notes, topic = _focus_of(student_id)
    if not notes:
        return "_⚠ Generate study material first._", None
    quiz = generate_quiz_json(notes, topic=topic, n=6)
    _mem(student_id)["quiz"] = quiz
    return quiz_to_table_md(quiz), quiz_file_clean(quiz)

def retry_plan(student_id, days, minutes):
    notes, topic = _focus_of(student_id)
    if not notes:
        return "_⚠ Generate study material first._"
    rows = generate_study_plan_rows(notes, topic, int(days or 7), int(minutes or 60))
    _mem(student_id)["plan_rows"] = rows
    return study_plan_to_md(rows, int(days or 7), int(minutes or 60))

def retry_flashcards(student_id):
    notes, topic = _focus_of(student_id)
    if not notes:
        return None
    return flashcards_pdf(generate_flashcards(notes, topic))

def retry_mindmap(student_id):
    notes, topic = _focus_of(student_id)
    if not notes:
        return "<p style='padding:12px'>⚠ Generate study material first.</p>"
    m = _mem(student_id)
    tree = generate_mindmap_tree(notes, topic, root_title=m.get("pdf_title", "Study Material"))
    m["tree"] = tree
    return build_mindmap_html(tree)

def retry_ppt(student_id):
    m = _mem(student_id)
    if not m or not m.get("focused_notes"):
        return None
    summary = m.get("summary") or layered_summary(m["focused_notes"], m.get("selected_topic", ""))
    return summary_to_pptx(summary, m.get("pdf_title", "Summary"))

def export_summary_pdf(student_id):
    m = _mem(student_id)
    if not m or not m.get("summary"):
        return None
    return summary_pdf(m["summary"], m.get("pdf_title", "Summary"))


# ============================================================
#  ACTIVE RECALL (RAG + fresh) & WEAK (dedup)
# ============================================================
def _qid(text):
    norm = re.sub(r"\s+", " ", (text or "").strip().lower())
    return hashlib.md5(norm.encode("utf-8")).hexdigest()[:10]

def start_practice(student_id):
    """Fresh recall quiz each click (uses RAG context + randomness)."""
    m = _mem(student_id)
    if not m or not m.get("focused_notes"):
        return ("<div style='padding:12px;color:#0f172a'>⚠ Generate study material "
                "first.</div>", gr.update(choices=[], value=None), "")
    prev = [q.get("question", "") for q in m.get("recall_quiz", [])]
    fresh = generate_fresh_recall_quiz(student_id, m["focused_notes"], prev,
                                       topic=m.get("selected_topic", ""))
    if not fresh:
        fresh = m.get("quiz", [])
    m["recall_quiz"] = fresh
    m["recall_index"] = 0
    return _show_recall(student_id, 0)

def _show_recall(student_id, index):
    m = _mem(student_id)
    quiz = (m or {}).get("recall_quiz", [])
    if not quiz:
        return ("<div style='padding:12px;color:#0f172a'>⚠ No recall questions "
                "available.</div>", gr.update(choices=[], value=None), "")
    index = index % len(quiz)
    m["recall_index"] = index
    q = quiz[index]
    opts = q.get("options", {})
    choices = [f"{k}) {opts[k]}" for k in ["A", "B", "C", "D"] if k in opts]
    qtype = q.get("question_type", "concept")
    sc = q.get("scenario_context", "")
    sc_html = ""
    if sc and str(sc).strip().lower() not in ("", "none", "empty"):
        sc_html = (f"<div style='background:#fdf2f8;border-left:4px solid #db2777;"
                   f"padding:10px 12px;border-radius:8px;margin:8px 0;color:#0f172a;"
                   f"font-style:italic;'>🎬 <b>Scenario:</b> {_esc(sc)}</div>")
    qhtml = (
        f"<div style='font-family:system-ui,Arial;border:1px solid #e2e8f0;"
        f"border-radius:14px;padding:16px;background:#fff;color:#0f172a;"
        f"box-shadow:0 2px 8px rgba(0,0,0,0.05);'>"
        f"<div style='display:flex;gap:8px;align-items:center;margin-bottom:8px;'>"
        f"<span style='background:#0f172a;color:#fff;padding:3px 12px;"
        f"border-radius:999px;font-size:12px;'>Q{index+1} / {len(quiz)}</span>"
        f"<span style='background:#ede9fe;color:#6d28d9;padding:3px 12px;"
        f"border-radius:999px;font-size:12px;'>{_esc(qtype)}</span></div>"
        f"{sc_html}"
        f"<div style='font-weight:700;font-size:15px;color:#0f172a;'>"
        f"{_esc(q.get('question',''))}</div></div>")
    return qhtml, gr.update(choices=choices, value=None), ""

def next_recall(student_id):
    m = _mem(student_id)
    idx = (m or {}).get("recall_index", 0)
    return _show_recall(student_id, idx + 1)

def submit_recall(student_id, selected):
    """Check answer, update progress, dedup weak questions. Returns HTML feedback."""
    sid = (student_id or "student123").strip()
    m = _mem(student_id)
    quiz = (m or {}).get("recall_quiz", [])
    if not quiz:
        return ("<div style='padding:10px;color:#0f172a'>⚠ Start practice first.</div>")
    if not selected:
        return ("<div style='padding:10px;color:#0f172a'>⚠ Please select an "
                "option.</div>")
    idx = m.get("recall_index", 0)
    q = quiz[idx]
    correct = str(q.get("correct_answer", q.get("answer", ""))).strip().upper()
    chosen = selected.split(")")[0].strip().upper()
    topic = q.get("topic", "General")

    data, rec = get_student_record(sid)
    rec["questions_attempted"] += 1
    if chosen == correct:
        rec["correct"] += 1
        msg = (f"<div style='padding:12px;border-radius:10px;background:#dcfce7;"
               f"border:1px solid #86efac;color:#0f172a;font-family:system-ui;'>"
               f"<b style='color:#16a34a;'>✅ Correct!</b><br>"
               f"💡 {_esc(q.get('explanation',''))}</div>")
    else:
        rec["wrong"] += 1
        rec["weak_topics"][topic] = rec["weak_topics"].get(topic, 0) + 1
        qid = _qid(q.get("question", ""))
        wq = rec["weak_questions"]
        if qid in wq:
            wq[qid]["mistakes"] += 1
            wq[qid]["last_attempt"] = datetime.now().strftime("%Y-%m-%d %H:%M")
        else:
            wq[qid] = {
                "question": q.get("question", ""), "answer": correct,
                "explanation": q.get("explanation", ""), "topic": topic,
                "mistakes": 1,
                "last_attempt": datetime.now().strftime("%Y-%m-%d %H:%M"),
            }
        msg = (f"<div style='padding:12px;border-radius:10px;background:#fee2e2;"
               f"border:1px solid #fecaca;color:#0f172a;font-family:system-ui;'>"
               f"<b style='color:#dc2626;'>❌ Wrong.</b> Correct answer: "
               f"<b>{_esc(correct)}</b><br>💡 {_esc(q.get('explanation',''))}</div>")
    rec["last_study"] = datetime.now().strftime("%Y-%m-%d %H:%M")
    data[sid] = rec
    save_progress(data)
    return msg

def review_weak(student_id):
    """Deduped weak questions, clean styled cards with BLACK text."""
    sid = (student_id or "student123").strip()
    _, rec = get_student_record(sid)
    wq = rec.get("weak_questions", {})
    if not wq:
        return ("<div style='padding:16px;color:#16a34a;font-family:system-ui;"
                "font-size:15px;font-weight:600'>🎉 No weak questions yet — "
                "great job!</div>")
    cards = ""
    for w in sorted(wq.values(), key=lambda x: -x.get("mistakes", 1)):
        cards += f"""
        <div style="border:1px solid #fecaca;border-left:5px solid #ef4444;
                    border-radius:12px;padding:14px;margin-bottom:12px;
                    background:#fff;color:#0f172a;">
          <div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:8px;
                      align-items:center;">
            <span style="background:#fee2e2;color:#b91c1c;padding:3px 12px;
                  border-radius:999px;font-size:12px;font-weight:700;">
                  ❌ {w.get('mistakes',1)} mistake(s)</span>
            <span style="background:#fef9c3;color:#854d0e;padding:3px 12px;
                  border-radius:999px;font-size:12px;">🏷️ {_esc(w.get('topic','General'))}</span>
            <span style="margin-left:auto;color:#0f172a;font-size:12px;">
                  🕒 {_esc(w.get('last_attempt','-'))}</span>
          </div>
          <div style="font-weight:700;color:#0f172a;margin-bottom:6px;">
            {_esc(w.get('question',''))}</div>
          <div style="color:#0f172a;margin-bottom:4px;">
            ✅ <b>Correct answer:</b> {_esc(w.get('answer',''))}</div>
          <div style="color:#0f172a;background:#f8fafc;padding:8px 10px;
               border-radius:8px;">💡 {_esc(w.get('explanation',''))}</div>
        </div>"""
    return (f"<div style='font-family:system-ui,Arial;color:#0f172a;'>"
            f"<h3 style='color:#0f172a;'>🔁 Weak Questions to Review</h3>{cards}</div>")


# ============================================================
#  Q&A + DASHBOARD HANDLERS
# ============================================================
def answer_question(question, student_id):
    if not question or not question.strip():
        return "❌ Please enter a question."
    return answer_from_notes(question, student_id)

def refresh_dashboard(student_id):
    return build_dashboard_html(student_id)

def download_quiz_pdf(student_id):
    """Make a styled PDF of the current main quiz."""
    m = _mem(student_id)
    if not m or not m.get("quiz"):
        return None
    return quiz_pdf_styled(m["quiz"], title=m.get("selected_topic") or m.get("pdf_title", "Quiz"))

def download_recall_pdf(student_id):
    """Make a styled PDF of the current active-recall practice set."""
    m = _mem(student_id)
    if not m or not m.get("recall_quiz"):
        return None
    return recall_pdf_styled(m["recall_quiz"],
                             title=m.get("selected_topic") or "Active Recall")


# ============================================================
#  GRADIO UI (STEP 12 prompt — tabs + sections, no API box)
# ============================================================
with gr.Blocks(title="📘 Study Buddy", theme=gr.themes.Soft()) as demo:

    # Hidden state: holds the logged-in user's email (used as student_id)
    sid = gr.State(value="")
    user_name = gr.State(value="")
    mins_in = gr.State(value=60)   # fixed 60 min/day (field removed from UI)

    # ========================================================
    #  LOGIN / SIGN-UP SCREEN  (shown first)
    # ========================================================
    with gr.Group(visible=True) as login_screen:
        gr.Markdown(
            "# 📘 Study Buddy — AI Learning Assistant\n"
            "Please **sign in** to continue. New here? Create an account in the "
            "**Sign Up** tab first."
        )
        with gr.Tab("🔑 Sign In"):
            gr.Markdown("Already have an account? Enter your email and password.")
            in_email = gr.Textbox(label="Email", placeholder="you@example.com")
            in_pw = gr.Textbox(label="Password", type="password")
            signin_btn = gr.Button("🔓 Sign In", variant="primary")
            signin_msg = gr.Markdown()
        with gr.Tab("📝 Sign Up"):
            gr.Markdown("New user? Create your account below.")
            up_name = gr.Textbox(label="Full Name", placeholder="Your name")
            up_email = gr.Textbox(label="Email", placeholder="you@example.com")
            up_pw = gr.Textbox(label="Password", type="password")
            up_confirm = gr.Textbox(label="Confirm Password", type="password")
            signup_btn = gr.Button("✅ Create Account", variant="primary")
            signup_msg = gr.Markdown()

    # ========================================================
    #  MAIN APP  (hidden until the user signs in)
    # ========================================================
    with gr.Group(visible=False) as main_app:
        with gr.Row():
            welcome = gr.Markdown("## 📘 Study Buddy")
            logout_btn = gr.Button("🚪 Log out", scale=0)
        gr.Markdown(
            "Upload a PDF (or type a topic) and get a focused summary, quiz, active "
            "recall, flashcards, mindmap, study plan and progress tracking. "
            "Q&A and Active Recall use RAG; everything else is topic-specific."
        )

        with gr.Tab("📄 Generate"):
            with gr.Accordion("1 · Upload / Topic", open=True):
                pdf = gr.File(label="Upload PDF(s)", file_types=[".pdf"],
                              file_count="multiple")
                with gr.Row():
                    topic = gr.Textbox(label="Enter Topic (optional — focuses everything)",
                                       placeholder="e.g. mise-en-scène")
                    days_in = gr.Number(label="Study days", value=7, precision=0)
                run = gr.Button("🚀 Generate Study Material", variant="primary")
            with gr.Accordion("2 · Topic Check / Status", open=True):
                status_box = gr.Textbox(label="Status", lines=2, interactive=False)

        with gr.Tab("📚 Learning"):
            with gr.Accordion("3 · Summary", open=True):
                summary = gr.Markdown()
                with gr.Row():
                    retry_sum_btn = gr.Button("🔄 Regenerate Summary")
                    export_sum_btn = gr.Button("📥 Download Summary PDF")
                summary_file = gr.File(label="Summary PDF")
            with gr.Accordion("4 · Study Plan (Table)", open=False):
                plan = gr.Markdown()
                with gr.Row():
                    retry_plan_btn = gr.Button("🔄 Regenerate Plan")
                    export_plan_btn = gr.Button("📥 Download Plan")
                plan_file = gr.File(label="Study Plan File")
            with gr.Accordion("9 · Interactive Mindmap (clickable)", open=False):
                mind = gr.HTML()
                retry_mind_btn = gr.Button("🔄 Regenerate Mindmap")

        with gr.Tab("📝 Quiz & Recall"):
            with gr.Accordion("5 · Quiz (Interactive Cards)", open=True):
                quiz_display = gr.HTML()
                with gr.Row():
                    retry_quiz_btn = gr.Button("🔄 Regenerate Quiz")
                    dl_quiz_btn = gr.Button("📥 Download Quiz PDF")
                quiz_file_out = gr.File(label="📥 Quiz PDF")
            with gr.Accordion("6 · Active Recall Practice (fresh each time)", open=False):
                gr.Markdown("Each session generates **new** scenario/application "
                            "questions (RAG-powered).")
                load_recall_btn = gr.Button("▶ Start New Practice", variant="primary")
                recall_q = gr.HTML()
                recall_opts = gr.Radio(label="Choose an answer", choices=[])
                with gr.Row():
                    submit_btn = gr.Button("✅ Submit", variant="primary")
                    next_btn = gr.Button("➡ Next Question")
                recall_feedback = gr.HTML()
                with gr.Row():
                    dl_recall_btn = gr.Button("📥 Download This Practice Set (PDF)")
                recall_file_out = gr.File(label="📥 Recall Practice PDF")
            with gr.Accordion("7 · Weak Questions (no duplicates)", open=False):
                weak_btn = gr.Button("Show Weak Questions")
                weak_box = gr.HTML()

        with gr.Tab("📥 Downloads & Progress"):
            with gr.Accordion("8 · Flashcards & All Downloads", open=True):
                gr.Markdown("Click any button to (re)generate that file, then "
                            "download it. Each file is styled by type.")
                with gr.Row():
                    flash = gr.File(label="🃏 Flashcards PDF")
                    ppt = gr.File(label="📊 PPT Slides")
                with gr.Row():
                    dl_summary2 = gr.File(label="📘 Summary PDF")
                    dl_plan2 = gr.File(label="🗓 Study Plan PDF")
                with gr.Row():
                    dl_quiz2 = gr.File(label="📝 Quiz PDF")
                    mind_file = gr.File(label="🧠 Mindmap HTML")
                with gr.Row():
                    progress_file = gr.File(label="📈 Progress Report (with charts)")
                gr.Markdown("**Generate / Refresh files:**")
                with gr.Row():
                    retry_flash_btn = gr.Button("🃏 Flashcards")
                    retry_ppt_btn = gr.Button("📊 PPT")
                    gen_summary2_btn = gr.Button("📘 Summary")
                with gr.Row():
                    gen_plan2_btn = gr.Button("🗓 Study Plan")
                    gen_quiz2_btn = gr.Button("📝 Quiz")
                    gen_progress_btn = gr.Button("📈 Progress Report")
            with gr.Accordion("10 · Progress Dashboard (graphical)", open=True):
                dashboard = gr.HTML()
                dash_btn = gr.Button("🔄 Refresh Dashboard")

        with gr.Tab("💬 Ask Questions"):
            gr.Markdown("Ask in **English, Roman Urdu, or Urdu**. Answers come from "
                        "your notes (RAG). Extra examples are clearly labelled.")
            text_question = gr.Textbox(
                label="Your question",
                placeholder="e.g. 'mise-en-scene kya hota hai?'  or in English / Urdu",
                lines=2)
            ask_btn = gr.Button("🤔 Get Answer", variant="primary")
            answer_output = gr.Textbox(label="Answer", lines=6, interactive=False)

    # ================= WIRING =================
    run.click(
        full_pipeline,
        [pdf, topic, sid, days_in, mins_in],
        [summary, quiz_display, plan, mind,
         flash, ppt, quiz_file_out, mind_file, progress_file,
         status_box, dashboard],
    )
    retry_sum_btn.click(retry_summary, sid, summary)
    export_sum_btn.click(export_summary_pdf, sid, summary_file)
    retry_plan_btn.click(retry_plan, [sid, days_in, mins_in], plan)
    export_plan_btn.click(
        lambda s, d, mn: study_plan_file(
            (_mem(s) or {}).get("plan_rows", []), int(d or 7), int(mn or 60)),
        [sid, days_in, mins_in], plan_file)
    retry_quiz_btn.click(retry_quiz, sid, [quiz_display, quiz_file_out])
    dl_quiz_btn.click(download_quiz_pdf, sid, quiz_file_out)
    retry_mind_btn.click(retry_mindmap, sid, mind)

    load_recall_btn.click(start_practice, sid,
                          [recall_q, recall_opts, recall_feedback])
    submit_btn.click(submit_recall, [sid, recall_opts], recall_feedback)
    next_btn.click(next_recall, sid, [recall_q, recall_opts, recall_feedback])
    dl_recall_btn.click(download_recall_pdf, sid, recall_file_out)
    weak_btn.click(review_weak, sid, weak_box)

    retry_flash_btn.click(retry_flashcards, sid, flash)
    retry_ppt_btn.click(retry_ppt, sid, ppt)
    # extra download buttons in section 8
    gen_summary2_btn.click(export_summary_pdf, sid, dl_summary2)
    gen_plan2_btn.click(
        lambda s, d, mn: study_plan_file(
            (_mem(s) or {}).get("plan_rows", []), int(d or 7), int(mn or 60)),
        [sid, days_in, mins_in], dl_plan2)
    gen_quiz2_btn.click(download_quiz_pdf, sid, dl_quiz2)
    gen_progress_btn.click(progress_report_file, sid, progress_file)
    dash_btn.click(refresh_dashboard, sid, dashboard)
    ask_btn.click(answer_question, [text_question, sid], answer_output)

    # ---------- LOGIN / SIGNUP / LOGOUT WIRING ----------
    def do_sign_in(email, password):
        """Validate login; if ok, hide login screen and show the app."""
        ok, msg, name = sign_in(email, password)
        if ok:
            email_clean = (email or "").strip().lower()
            return (
                msg,                                   # signin_msg
                gr.update(visible=False),              # hide login_screen
                gr.update(visible=True),               # show main_app
                email_clean,                           # sid state = email
                name,                                  # user_name state
                f"## 👋 Welcome, {name}!",             # welcome
                build_dashboard_html(email_clean),     # preload dashboard
            )
        # failed: stay on login screen
        return (msg, gr.update(visible=True), gr.update(visible=False),
                "", "", "## 📘 Study Buddy", "")

    def do_sign_up(name, email, pw, confirm):
        """Create account; on success, tell the user to sign in."""
        ok, msg = sign_up(email, name, pw, confirm)
        return msg

    def do_logout():
        """Return to the login screen and clear the session id."""
        return (gr.update(visible=True),   # show login_screen
                gr.update(visible=False),  # hide main_app
                "", "",                    # clear sid + user_name
                "")                        # clear signin_msg

    signin_btn.click(
        do_sign_in, [in_email, in_pw],
        [signin_msg, login_screen, main_app, sid, user_name, welcome, dashboard])
    signup_btn.click(do_sign_up, [up_name, up_email, up_pw, up_confirm], signup_msg)
    logout_btn.click(do_logout, None,
                     [login_screen, main_app, sid, user_name, signin_msg])

    gr.Markdown(
        "> 📝 **How to use:** put your OpenAI key in the `API_KEY` variable at the "
        "top → run → **Sign Up** (first time) → **Sign In** → upload a PDF "
        "(optionally type a topic) → **Generate Study Material**."
    )

# share=True gives a public link (needed in Colab)
demo.launch(share=True)

